"""PPTX parser implementation producing slide and note text blocks."""

from __future__ import annotations

import io
import zipfile
from collections.abc import Iterable, Mapping
from typing import Any, List, Optional, Set, Tuple, Dict

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.exc import PackageNotFoundError
from pptx.opc.constants import RELATIONSHIP_TYPE

from documents.contract_utils import (
    is_bcp47_like,
    normalize_media_type,
    normalize_string,
    truncate_text,
)
from documents.parsers import (
    DocumentParser,
    ParsedAsset,
    ParsedResult,
    ParsedTextBlock,
    build_parsed_asset_with_meta,
    build_parsed_result,
    build_parsed_text_block,
)

_PPTX_MEDIA_TYPES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

_PPTX_MAX_ARCHIVE_ENTRIES = 10_000
_PPTX_MAX_UNCOMPRESSED_SIZE = 500 * 1024 * 1024
_PPTX_MAX_ASSET_SIZE = 25 * 1024 * 1024
_PPTX_MAX_TOTAL_ASSET_SIZE = 200 * 1024 * 1024

_IMAGE_EXTENSION_FALLBACK = {
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".png": "image/png",
    ".gif": "image/gif",
    ".bmp": "image/bmp",
    ".tif": "image/tiff",
    ".tiff": "image/tiff",
}


def _normalise_media_type(value: Optional[str]) -> Optional[str]:
    if not value:
        return None
    try:
        return normalize_media_type(value)
    except ValueError:
        return None


def _extract_media_type(document: Any) -> Optional[str]:
    media_type = getattr(document, "media_type", None)
    if not media_type and isinstance(document, Mapping):
        media_type = document.get("media_type")
    return _normalise_media_type(media_type)


def _extract_blob(document: Any) -> Any:
    blob = getattr(document, "blob", None)
    if blob is None and isinstance(document, Mapping):
        blob = document.get("blob")
    return blob


def _extract_blob_media_type(blob: Any) -> Optional[str]:
    media_type = getattr(blob, "media_type", None)
    if media_type is None and isinstance(blob, Mapping):
        media_type = blob.get("media_type")
    return _normalise_media_type(media_type)


def _blob_payload(blob: Any) -> Optional[bytes]:
    if blob is None:
        return None
    if hasattr(blob, "decoded_payload"):
        payload = blob.decoded_payload()
        if isinstance(payload, bytes):
            return payload
    data = getattr(blob, "content", None)
    if isinstance(data, (bytes, bytearray)):
        return bytes(data)
    if isinstance(blob, Mapping):
        inline = blob.get("content")
        if isinstance(inline, (bytes, bytearray)):
            return bytes(inline)
        base64_value = blob.get("base64")
        if isinstance(base64_value, (bytes, bytearray)):
            return bytes(base64_value)
    if isinstance(blob, (bytes, bytearray)):
        return bytes(blob)
    return None


def _validate_archive(archive: zipfile.ZipFile) -> None:
    infos = archive.infolist()
    if len(infos) > _PPTX_MAX_ARCHIVE_ENTRIES:
        raise ValueError("pptx_archive_limits")
    total_size = 0
    for info in infos:
        total_size += max(info.file_size, 0)
        if total_size > _PPTX_MAX_UNCOMPRESSED_SIZE:
            raise ValueError("pptx_archive_limits")


def _looks_like_pptx(payload: Optional[bytes]) -> bool:
    if not payload:
        return False
    if not payload.startswith(b"PK"):
        return False
    try:
        with zipfile.ZipFile(io.BytesIO(payload)) as archive:
            try:
                _validate_archive(archive)
            except ValueError:
                return False
            names = set(archive.namelist())
            return "ppt/presentation.xml" in names
    except zipfile.BadZipFile:
        return False


def _iter_shapes(shapes: Iterable[Any]) -> Iterable[Any]:
    for shape in shapes:
        yield shape
        shape_type = None
        try:
            shape_type = getattr(shape, "shape_type", None)
        except NotImplementedError:  # pragma: no cover - defensive
            shape_type = None
        if shape_type == MSO_SHAPE_TYPE.GROUP and hasattr(shape, "shapes"):
            yield from _iter_shapes(shape.shapes)


def _collect_shape_text(shape: Any) -> Tuple[List[str], Set[str]]:
    paragraphs: List[str] = []
    languages: Set[str] = set()
    if getattr(shape, "has_text_frame", False):
        try:
            text_frame = shape.text_frame
        except AttributeError:  # pragma: no cover - defensive
            text_frame = None
        if text_frame is not None:
            for paragraph in text_frame.paragraphs:
                text = normalize_string(paragraph.text)
                if text:
                    paragraphs.append(text)
        try:
            for raw_lang in shape.element.xpath(".//a:rPr/@lang"):
                candidate = normalize_string(raw_lang)
                if candidate and is_bcp47_like(candidate):
                    languages.add(candidate)
        except AttributeError:  # pragma: no cover - defensive
            pass
    return paragraphs, languages


def _collect_slide_content(slide: Any) -> Tuple[str, Optional[str], Optional[str]]:
    paragraphs: List[str] = []
    languages: Set[str] = set()
    for shape in _iter_shapes(slide.shapes):
        parts, langs = _collect_shape_text(shape)
        if parts:
            paragraphs.extend(parts)
        if langs:
            languages.update(langs)
    slide_text = "\n".join(paragraphs).strip()

    title_text: Optional[str] = None
    try:
        title_shape = slide.shapes.title
    except AttributeError:  # pragma: no cover - defensive
        title_shape = None
    if title_shape is not None and getattr(title_shape, "text", None):
        title_candidate = normalize_string(title_shape.text)
        if title_candidate:
            title_text = title_candidate

    slide_language = next(iter(languages)) if len(languages) == 1 else None
    return slide_text, title_text, slide_language


def _collect_notes_content(notes_slide: Any) -> Tuple[str, Set[str]]:
    paragraphs: List[str] = []
    languages: Set[str] = set()
    if notes_slide is None:
        return "", languages
    for shape in _iter_shapes(notes_slide.shapes):
        parts, langs = _collect_shape_text(shape)
        if parts:
            paragraphs.extend(parts)
        if langs:
            languages.update(langs)
    text = "\n".join(paragraphs).strip()
    return text, languages


def _infer_media_type(image: Any) -> str:
    media_type = _normalise_media_type(getattr(image, "content_type", None))
    if media_type:
        return media_type
    extension = getattr(image, "ext", None)
    if extension:
        fallback = _IMAGE_EXTENSION_FALLBACK.get(f".{extension.lower()}")
        media_type = _normalise_media_type(fallback)
        if media_type:
            return media_type
    return "image/unknown"


def _extract_alt_text(shape: Any) -> Optional[str]:
    raw_alt = getattr(shape, "alternative_text", None)
    if isinstance(raw_alt, str) and raw_alt:
        candidate = normalize_string(raw_alt)
        if candidate:
            return candidate
    pic = getattr(shape, "_pic", None)
    c_nv_pr = getattr(getattr(pic, "nvPicPr", None), "cNvPr", None)
    if c_nv_pr is not None:
        for attr in ("descr", "title"):
            value = c_nv_pr.get(attr)
            if value:
                candidate = normalize_string(value)
                if candidate:
                    return candidate
    return None


class PptxDocumentParser(DocumentParser):
    """Parse PPTX payloads into slide, note, and asset results."""

    def can_handle(self, document: Any) -> bool:
        media_type = _extract_media_type(document)
        if media_type in _PPTX_MEDIA_TYPES:
            return True
        blob = _extract_blob(document)
        blob_media_type = _extract_blob_media_type(blob)
        if blob_media_type in _PPTX_MEDIA_TYPES:
            return True
        payload = _blob_payload(blob)
        return _looks_like_pptx(payload)

    def parse(self, document: Any, config: Any) -> ParsedResult:
        blob = _extract_blob(document)
        payload = _blob_payload(blob)
        if not payload:
            raise ValueError("pptx_blob_missing")

        notes_enabled = bool(getattr(config, "enable_notes_in_pptx", False))
        emit_empty_slides = bool(getattr(config, "emit_empty_slides", True))

        buffer = io.BytesIO(payload)
        try:
            with zipfile.ZipFile(buffer) as archive:
                _validate_archive(archive)
                if "ppt/presentation.xml" not in archive.namelist():
                    raise ValueError("pptx_presentation_missing")
        except zipfile.BadZipFile as exc:  # pragma: no cover - defensive guard
            raise ValueError("pptx_invalid_package") from exc
        finally:
            buffer.seek(0)

        try:
            presentation = Presentation(buffer)
        except PackageNotFoundError as exc:  # pragma: no cover - defensive guard
            raise ValueError("pptx_invalid_package") from exc
        except KeyError as exc:  # pragma: no cover - defensive guard
            raise ValueError("pptx_presentation_missing") from exc

        slide_blocks: List[ParsedTextBlock] = []
        note_blocks: List[ParsedTextBlock] = []
        assets: List[ParsedAsset] = []
        images_count = 0
        assets_total_size = 0

        for index, slide in enumerate(presentation.slides):
            slide_text, title, slide_language = _collect_slide_content(slide)
            section_segments = ["Slide", str(index + 1)]
            if title:
                section_segments.append(title)
            section_path = tuple(section_segments)

            if slide_text or emit_empty_slides:
                block_text = slide_text or f"Slide {index + 1}"
                slide_blocks.append(
                    build_parsed_text_block(
                        text=block_text,
                        kind="slide",
                        section_path=section_path,
                        page_index=index,
                        language=slide_language,
                    )
                )

            slide_context = slide_text or None
            slide_context_truncated = (
                truncate_text(slide_context, 512) if slide_context else None
            )

            combined_notes = ""
            note_languages: Set[str] = set()
            if notes_enabled:
                notes_texts: List[str] = []
                for rel in slide.part.rels.values():
                    if rel.reltype != RELATIONSHIP_TYPE.NOTES_SLIDE:
                        continue
                    notes_part = getattr(rel, "target_part", None)
                    notes_slide = getattr(notes_part, "notes_slide", None)
                    if notes_slide is None:
                        continue
                    note_text, languages = _collect_notes_content(notes_slide)
                    if note_text:
                        notes_texts.append(note_text)
                    if languages:
                        note_languages.update(languages)
                combined_notes = "\n\n".join(notes_texts).strip()

                if combined_notes:
                    if note_languages:
                        notes_language = (
                            next(iter(note_languages))
                            if len(note_languages) == 1
                            else None
                        )
                    else:
                        notes_language = slide_language
                    note_blocks.append(
                        build_parsed_text_block(
                            text=combined_notes,
                            kind="note",
                            section_path=section_path,
                            page_index=index,
                            language=notes_language,
                        )
                    )

            notes_context = (
                truncate_text(combined_notes, 512) if combined_notes else None
            )

            parent_ref = f"slide:{index + 1}"
            asset_seq = 0

            for shape in _iter_shapes(slide.shapes):
                shape_type = getattr(shape, "shape_type", None)
                if shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                    except (
                        AttributeError,
                        KeyError,
                    ) as exc:  # pragma: no cover - defensive guard
                        raise ValueError("pptx_image_missing") from exc
                    content = bytes(image.blob)
                    if len(content) > _PPTX_MAX_ASSET_SIZE:
                        raise ValueError("pptx_asset_too_large")
                    if assets_total_size + len(content) > _PPTX_MAX_TOTAL_ASSET_SIZE:
                        raise ValueError("pptx_assets_total_too_large")
                    assets_total_size += len(content)
                    media_type = _infer_media_type(image)
                    alt_text = _extract_alt_text(shape)

                    context_before = slide_context_truncated
                    context_after = notes_context or slide_context_truncated
                    if context_after is None and alt_text:
                        context_after = truncate_text(alt_text, 512)
                    asset_seq += 1
                    shape_id = getattr(shape, "shape_id", None)
                    if shape_id is not None:
                        locator = f"{parent_ref}:shape:{shape_id}"
                    else:
                        locator = f"{parent_ref}:image:{asset_seq}"
                    candidates: List[Tuple[str, str]] = []
                    if alt_text:
                        candidates.append(("alt_text", alt_text))
                    if notes_context:
                        candidates.append(("notes", notes_context))
                    if slide_context_truncated:
                        candidates.append(("context_after", slide_context_truncated))
                    if context_before:
                        candidates.append(("context_before", context_before))
                    metadata: Dict[str, Any] = {
                        "parent_ref": parent_ref,
                        "locator": locator,
                    }
                    if candidates:
                        metadata["caption_candidates"] = candidates

                    assets.append(
                        build_parsed_asset_with_meta(
                            media_type=media_type,
                            content=content,
                            page_index=index,
                            context_before=context_before,
                            context_after=context_after,
                            metadata=metadata,
                        )
                    )
                    images_count += 1
                    continue

                if getattr(shape, "has_chart", False):
                    chart = getattr(shape, "chart", None)
                    if chart is None:
                        continue
                    chart_part = getattr(chart, "part", None)
                    if chart_part is None:
                        continue
                    partname = getattr(chart_part, "partname", None)
                    file_uri = str(partname) if partname is not None else None
                    if not file_uri:
                        continue
                    chart_locator = f"{parent_ref}:chart:{getattr(shape, 'shape_id', asset_seq + 1)}"
                    chart_media_type = getattr(
                        chart_part,
                        "content_type",
                        "application/vnd.openxmlformats-officedocument.drawingml.chart+xml",
                    )
                    candidates: List[Tuple[str, str]] = []
                    if notes_context:
                        candidates.append(("notes", notes_context))
                    if slide_context_truncated:
                        candidates.append(("context_after", slide_context_truncated))
                    metadata: Dict[str, Any] = {
                        "parent_ref": parent_ref,
                        "locator": chart_locator,
                        "asset_kind": "chart_source",
                    }
                    if candidates:
                        metadata["caption_candidates"] = candidates

                    assets.append(
                        build_parsed_asset_with_meta(
                            media_type=chart_media_type,
                            file_uri=file_uri,
                            page_index=index,
                            context_before=slide_context_truncated,
                            context_after=notes_context or slide_context_truncated,
                            metadata=metadata,
                        )
                    )

        statistics = {
            "parser.slides": len(slide_blocks),
            "parser.notes": len(note_blocks),
            "assets.images": images_count,
        }

        return build_parsed_result(
            text_blocks=tuple(slide_blocks + note_blocks),
            assets=tuple(assets),
            statistics=statistics,
        )
