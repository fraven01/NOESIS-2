import base64
import hashlib
import io
import zipfile
from xml.etree import ElementTree as ET
from xml.sax.saxutils import escape

import pytest
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.util import Inches

import documents.parsers_pptx as pptx_module

from documents.contracts import InlineBlob
from documents.parsers_pptx import PptxDocumentParser
from documents.pipeline import DocumentPipelineConfig

_PNG_BASE64 = "iVBORw0KGgoAAAANSUhEUgAAAAgAAAAICAIAAABLbSncAAAAFElEQVR4nGOsfq7LgA0wYRUdtBIAO/8Bn565mEEAAAAASUVORK5CYII="
_PNG_BYTES = base64.b64decode(_PNG_BASE64)

_NOTES_RELATIONSHIP = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide"
)


def _inline_blob_from_payload(payload: bytes, media_type: str) -> InlineBlob:
    return InlineBlob(
        type="inline",
        media_type=media_type,
        base64=base64.b64encode(payload).decode("ascii"),
        sha256=hashlib.sha256(payload).hexdigest(),
        size=len(payload),
    )


def _make_document(payload: bytes, media_type: str | None = None) -> dict:
    effective_type = (
        media_type
        or "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    return {"blob": _inline_blob_from_payload(payload, effective_type)}


def _set_language(text_frame, language) -> None:
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.language_id = language


def _build_presentation_bytes(builder) -> bytes:
    presentation = Presentation()
    builder(presentation)
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _sample_presentation_payload() -> bytes:
    def builder(presentation: Presentation) -> None:
        slide1 = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide1.shapes.title.text = "Quarter Update"
        _set_language(slide1.shapes.title.text_frame, MSO_LANGUAGE_ID.ENGLISH_US)
        body1 = slide1.placeholders[1]
        body1.text = "Highlights\nRevenue up 15%"
        _set_language(body1.text_frame, MSO_LANGUAGE_ID.ENGLISH_US)
        notes1 = slide1.notes_slide
        notes1.notes_text_frame.text = "Discuss quarterly growth narrative"
        _set_language(notes1.notes_text_frame, MSO_LANGUAGE_ID.ENGLISH_US)
        image_stream = io.BytesIO(_PNG_BYTES)
        image_stream.seek(0)
        picture1 = slide1.shapes.add_picture(image_stream, Inches(0.2), Inches(0.2))
        picture1._pic.nvPicPr.cNvPr.set("descr", "Chart overview")

        slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide2.shapes.title.text = "Next Steps"
        _set_language(slide2.shapes.title.text_frame, MSO_LANGUAGE_ID.ENGLISH_US)
        body2 = slide2.placeholders[1]
        body2.text = "Finalize roadmap"
        _set_language(body2.text_frame, MSO_LANGUAGE_ID.ENGLISH_US)
        notes2 = slide2.notes_slide
        notes2.notes_text_frame.text = "Follow up with stakeholders"
        _set_language(notes2.notes_text_frame, MSO_LANGUAGE_ID.ENGLISH_US)
        image_stream2 = io.BytesIO(_PNG_BYTES)
        image_stream2.seek(0)
        picture2 = slide2.shapes.add_picture(image_stream2, Inches(1), Inches(1))
        picture2._pic.nvPicPr.cNvPr.set("descr", "Roadmap illustration")

    return _build_presentation_bytes(builder)


def _add_additional_note(payload: bytes, text: str) -> bytes:
    parts: dict[str, bytes] = {}
    with zipfile.ZipFile(io.BytesIO(payload)) as archive:
        for name in archive.namelist():
            parts[name] = archive.read(name)

    new_part_name = "ppt/notesSlides/notesSlideExtra.xml"
    counter = 1
    while new_part_name in parts:
        counter += 1
        new_part_name = f"ppt/notesSlides/notesSlideExtra{counter}.xml"

    escaped_text = escape(text)
    note_xml = f"""<?xml version='1.0' encoding='UTF-8' standalone='yes'?>
<p:notes xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr>
        <p:cNvPr id="1" name=""/>
        <p:cNvGrpSpPr/>
        <p:nvPr/>
      </p:nvGrpSpPr>
      <p:grpSpPr/>
      <p:sp>
        <p:nvSpPr>
          <p:cNvPr id="2" name="Notes Placeholder"/>
          <p:cNvSpPr/>
          <p:nvPr/>
        </p:nvSpPr>
        <p:spPr/>
        <p:txBody>
          <a:bodyPr/>
          <a:lstStyle/>
          <a:p>
            <a:r>
              <a:t>{escaped_text}</a:t>
            </a:r>
          </a:p>
        </p:txBody>
      </p:sp>
    </p:spTree>
  </p:cSld>
</p:notes>
""".encode(
        "utf-8"
    )
    parts[new_part_name] = note_xml

    content_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
    override_tag = f"{{{content_ns}}}Override"
    content_root = ET.fromstring(parts["[Content_Types].xml"])
    part_name_attr = f"/{new_part_name}"
    if not any(
        elem.get("PartName") == part_name_attr
        for elem in content_root.findall(override_tag)
    ):
        ET.SubElement(
            content_root,
            override_tag,
            PartName=part_name_attr,
            ContentType="application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml",
        )
    parts["[Content_Types].xml"] = ET.tostring(
        content_root, encoding="utf-8", xml_declaration=True
    )

    rels_path = "ppt/slides/_rels/slide1.xml.rels"
    rels_root = ET.fromstring(parts[rels_path])
    rel_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
    rel_tag = f"{{{rel_ns}}}Relationship"
    existing_ids = {elem.get("Id") for elem in rels_root.findall(rel_tag)}
    counter = 1
    while f"rId{counter}" in existing_ids:
        counter += 1
    ET.SubElement(
        rels_root,
        rel_tag,
        Id=f"rId{counter}",
        Type=_NOTES_RELATIONSHIP,
        Target=f"../notesSlides/{new_part_name.split('/')[-1]}",
    )
    parts[rels_path] = ET.tostring(rels_root, encoding="utf-8", xml_declaration=True)

    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
        for name, content in parts.items():
            archive.writestr(name, content)
    return buffer.getvalue()


def test_pptx_parser_extracts_slides_notes_and_assets():
    parser = PptxDocumentParser()
    payload = _sample_presentation_payload()
    document = _make_document(payload)

    result = parser.parse(document, DocumentPipelineConfig())

    kinds = [block.kind for block in result.text_blocks]
    assert kinds == ["slide", "slide", "note", "note"]

    slide_one = result.text_blocks[0]
    assert slide_one.text == "Quarter Update\nHighlights\nRevenue up 15%"
    assert slide_one.section_path == ("Slide", "1", "Quarter Update")
    assert slide_one.language == "en-US"

    note_one = result.text_blocks[2]
    assert note_one.text == "Discuss quarterly growth narrative"
    assert note_one.section_path == ("Slide", "1", "Quarter Update")
    assert note_one.language == "en-US"

    assert len(result.assets) == 2
    asset = result.assets[0]
    assert asset.media_type == "image/png"
    assert asset.context_before.startswith("Quarter Update")
    assert asset.context_after.startswith("Discuss quarterly")
    assert asset.metadata.get("parent_ref") == "slide:1"
    assert asset.metadata.get("locator", "").startswith("slide:1:")
    candidates = asset.metadata.get("caption_candidates", [])
    assert candidates and candidates[0] == ("alt_text", "Chart overview")
    second_asset = result.assets[1]
    assert second_asset.metadata.get("parent_ref") == "slide:2"
    second_candidates = second_asset.metadata.get("caption_candidates", [])
    assert second_candidates and second_candidates[0] == (
        "alt_text",
        "Roadmap illustration",
    )
    assert result.statistics["parser.slides"] == 2
    assert result.statistics["parser.notes"] == 2
    assert result.statistics["assets.images"] == 2
    assert result.statistics["parse.blocks.total"] == 4
    assert result.statistics["parse.assets.total"] == 2


def test_pptx_parser_respects_notes_flag():
    parser = PptxDocumentParser()
    payload = _sample_presentation_payload()
    document = _make_document(payload)

    result = parser.parse(document, DocumentPipelineConfig(enable_notes_in_pptx=False))

    assert [block.kind for block in result.text_blocks] == ["slide", "slide"]
    assert result.statistics["parser.notes"] == 0


def test_pptx_parser_concatenates_multiple_notes():
    parser = PptxDocumentParser()
    payload = _add_additional_note(
        _sample_presentation_payload(), "Second note payload"
    )
    document = _make_document(payload)

    result = parser.parse(document, DocumentPipelineConfig())

    note_texts = [
        block.text
        for block in result.text_blocks
        if block.kind == "note" and block.page_index == 0
    ]
    assert note_texts == ["Discuss quarterly growth narrative\n\nSecond note payload"]


def test_pptx_parser_empty_slide_handling():
    parser = PptxDocumentParser()

    def builder(presentation: Presentation) -> None:
        slide1 = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide1.shapes.title.text = "With Text"
        _set_language(slide1.shapes.title.text_frame, MSO_LANGUAGE_ID.ENGLISH_US)
        body1 = slide1.placeholders[1]
        body1.text = "Body"
        _set_language(body1.text_frame, MSO_LANGUAGE_ID.ENGLISH_US)

        slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide2.shapes.title.text = ""
        slide2.placeholders[1].text = ""

    payload = _build_presentation_bytes(builder)
    document = _make_document(payload)

    default_result = parser.parse(document, DocumentPipelineConfig())
    assert [block.text for block in default_result.text_blocks] == [
        "With Text\nBody",
        "Slide 2",
    ]

    skip_result = parser.parse(
        document, DocumentPipelineConfig(emit_empty_slides=False)
    )
    assert [block.text for block in skip_result.text_blocks] == ["With Text\nBody"]


def test_pptx_parser_enforces_total_asset_budget(monkeypatch):
    parser = PptxDocumentParser()
    payload = _sample_presentation_payload()
    document = _make_document(payload)
    config = DocumentPipelineConfig()

    monkeypatch.setattr(pptx_module, "_PPTX_MAX_TOTAL_ASSET_SIZE", 100)

    with pytest.raises(ValueError) as exc:
        parser.parse(document, config)

    assert str(exc.value) == "pptx_assets_total_too_large"


def test_pptx_parser_alt_text_context_when_notes_disabled():
    parser = PptxDocumentParser()

    def builder(presentation: Presentation) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        image_stream = io.BytesIO(_PNG_BYTES)
        image_stream.seek(0)
        picture = slide.shapes.add_picture(image_stream, Inches(1), Inches(1))
        picture._pic.nvPicPr.cNvPr.set("descr", "Standalone diagram")

    payload = _build_presentation_bytes(builder)
    document = _make_document(payload)

    result = parser.parse(document, DocumentPipelineConfig(enable_notes_in_pptx=False))

    assert len(result.assets) == 1
    asset = result.assets[0]
    assert asset.context_before is None
    assert asset.context_after == "Standalone diagram"
    alt_candidates = asset.metadata.get("caption_candidates", [])
    assert alt_candidates and alt_candidates[0] == ("alt_text", "Standalone diagram")


def test_pptx_parser_emits_chart_source_assets():
    parser = PptxDocumentParser()

    def builder(presentation: Presentation) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "Chart Slide"
        chart_data = CategoryChartData()
        chart_data.categories = ["Q1", "Q2"]
        chart_data.add_series("Series", (1, 2))
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(3),
            chart_data,
        )

    payload = _build_presentation_bytes(builder)
    document = _make_document(payload)

    result = parser.parse(document, DocumentPipelineConfig())

    chart_assets = [
        asset
        for asset in result.assets
        if asset.metadata.get("asset_kind") == "chart_source"
    ]
    assert chart_assets
    chart_asset = chart_assets[0]
    assert chart_asset.metadata.get("parent_ref") == "slide:1"
    assert chart_asset.metadata.get("locator", "").startswith("slide:1:chart:")
    assert chart_asset.file_uri and chart_asset.file_uri.endswith(".xml")


def test_pptx_parser_language_with_mixed_runs_is_none():
    parser = PptxDocumentParser()

    def builder(presentation: Presentation) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        title_frame = slide.shapes.title.text_frame
        paragraph = title_frame.paragraphs[0]
        paragraph.clear()
        run_en = paragraph.add_run()
        run_en.text = "Hello"
        run_en.font.language_id = MSO_LANGUAGE_ID.ENGLISH_US
        run_fr = paragraph.add_run()
        run_fr.text = "Bonjour"
        run_fr.font.language_id = MSO_LANGUAGE_ID.FRENCH

    payload = _build_presentation_bytes(builder)
    document = _make_document(payload)

    result = parser.parse(document, DocumentPipelineConfig())

    assert result.text_blocks[0].language is None


def test_pptx_parser_note_language_falls_back_to_notes_metadata():
    parser = PptxDocumentParser()

    def builder(presentation: Presentation) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "No Language"
        slide.placeholders[1].text = "Body without tags"
        notes = slide.notes_slide
        notes.notes_text_frame.text = "Hinweis zum Ablauf"
        _set_language(notes.notes_text_frame, MSO_LANGUAGE_ID.GERMAN)

    payload = _build_presentation_bytes(builder)
    document = _make_document(payload)

    result = parser.parse(document, DocumentPipelineConfig())

    slide_block = result.text_blocks[0]
    assert slide_block.language is None

    note_block = result.text_blocks[1]
    assert note_block.language == "de-DE"


def test_pptx_parser_can_handle_payload_without_media_type():
    parser = PptxDocumentParser()
    payload = _sample_presentation_payload()
    blob = _inline_blob_from_payload(payload, "application/octet-stream")

    assert parser.can_handle({"blob": blob}) is True
